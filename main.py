from fastapi import FastAPI, UploadFile, File, Form, Query, Request, BackgroundTasks
from fastapi.responses import FileResponse, HTMLResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from typing import List, Optional
import os
import shutil
import uuid, zipfile
import asyncio 
from PIL import Image
from tools import router as tools_router

# Import your backend modules
from pypdf import PdfReader, PdfWriter
from pdf2docx import Converter
import fitz  # PyMuPDF for text extraction & image generation
from gtts import gTTS  # Text to speech

app = FastAPI()

# Folders
os.makedirs("uploads", exist_ok=True)
os.makedirs("output", exist_ok=True)

# Static & templates
app.mount("/Static", StaticFiles(directory="Static"), name="Static")
templates = Jinja2Templates(directory="templates")

UPLOAD_DIR = "uploads"
OUTPUT_DIR = "output"

# ------------------ PAGES ------------------
@app.get("/", response_class=HTMLResponse)
async def home(request: Request):
    return templates.TemplateResponse("Home.html", {"request": request})

@app.get("/pdf_Convertion", response_class=HTMLResponse)
async def pdf_convert(request: Request):
    return templates.TemplateResponse("PDF_Convert.html", {"request": request})

@app.get("/Image_Convertion", response_class=HTMLResponse)
async def image_page(request: Request):
    return templates.TemplateResponse("Image_Convert.html", {"request": request})

@app.get("/resize-image", response_class=HTMLResponse)
async def resize_page(request: Request):
    return templates.TemplateResponse("resize-editor.html", {"request": request})

@app.get("/rotate-image", response_class=HTMLResponse)
async def rotate_page(request: Request):
    return templates.TemplateResponse("rotate-editor.html", {"request": request})

# Function to delay deletion
async def delete_file_delayed(path: str, delay_seconds: int = 300):
    """Wait before deleting to allow download managers to grab all file chunks"""
    await asyncio.sleep(delay_seconds)
    try:
        if os.path.exists(path):
            os.remove(path)
    except Exception:
        pass

# ------------------ PDF CONVERSION FUNCTIONS ------------------
def pdf_to_word_converter(input_path: str) -> str:
    output_path = os.path.join(OUTPUT_DIR, f"converted_{uuid.uuid4()}.docx")
    cv = Converter(input_path)
    cv.convert(output_path)
    cv.close()
    return output_path

def pdf_to_text_converter(input_path: str) -> str:
    output_path = os.path.join(OUTPUT_DIR, f"converted_{uuid.uuid4()}.txt")
    doc = fitz.open(input_path)
    text = "".join([page.get_text() for page in doc])
    doc.close()
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(text)
    return output_path

def pdf_to_mp3_converter(input_path: str) -> str:
    doc = fitz.open(input_path)
    text = "".join([page.get_text() for page in doc])
    doc.close()
    output_path = os.path.join(OUTPUT_DIR, f"converted_{uuid.uuid4()}.mp3")
    tts = gTTS(text=text, lang='en', slow=False)
    tts.save(output_path)
    return output_path

def pdf_to_excel_converter(input_path: str) -> str:
    output_path = os.path.join(OUTPUT_DIR, f"converted_{uuid.uuid4()}.xlsx")
    import openpyxl
    doc = fitz.open(input_path)
    wb = openpyxl.Workbook()
    ws = wb.active
    row = 1
    for page in doc:
        text = page.get_text()
        for line in text.split('\n'):
            ws.cell(row=row, column=1, value=line)
            row += 1
    doc.close()
    wb.save(output_path)
    return output_path

def pdf_to_ppt_converter(input_path: str) -> str:
    from pptx import Presentation
    from pptx.util import Inches
    output_path = os.path.join(OUTPUT_DIR, f"converted_{uuid.uuid4()}.pptx")
    doc = fitz.open(input_path)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    for i in range(len(doc)):
        page = doc.load_page(i)
        pix = page.get_pixmap(dpi=150)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        temp_img_path = os.path.join(OUTPUT_DIR, f"temp_slide_{i}.png")
        pix.save(temp_img_path)
        slide.shapes.add_picture(temp_img_path, 0, 0, width=prs.slide_width)
        os.remove(temp_img_path)
    doc.close()
    prs.save(output_path)
    return output_path

def compress_pdf_converter(input_path: str) -> str:
    output_path = os.path.join(OUTPUT_DIR, f"compressed_{uuid.uuid4()}.pdf")
    reader = PdfReader(input_path)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    for page in writer.pages:
        page.compress_content_streams() 
    writer.add_metadata(reader.metadata) 
    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path

def merge_pdf_converter(input_paths: list) -> str:
    output_path = os.path.join(OUTPUT_DIR, f"merged_{uuid.uuid4()}.pdf")
    writer = PdfWriter()
    for pdf_path in input_paths:
        reader = PdfReader(pdf_path)
        for page in reader.pages:
            writer.add_page(page)
    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path

def pdf_to_image_converter(input_path: str) -> str:
    unique_id = str(uuid.uuid4())
    job_folder = os.path.join(OUTPUT_DIR, f"img_job_{unique_id}")
    os.makedirs(job_folder, exist_ok=True)
    try:
        doc = fitz.open(input_path)
        if len(doc) == 0: raise Exception("This PDF has no pages.")
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(dpi=150)
            pix.save(os.path.join(job_folder, f"page_{page_num + 1}.png"))
        doc.close()
        zip_path = os.path.join(OUTPUT_DIR, f"images_{unique_id}.zip")
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for root, dirs, files in os.walk(job_folder):
                for file in files: zipf.write(os.path.join(root, file), arcname=file)
        return zip_path
    finally:
        shutil.rmtree(job_folder, ignore_errors=True)

def split_pdf_converter(input_path: str) -> str:
    unique_id = str(uuid.uuid4())
    job_folder = os.path.join(OUTPUT_DIR, f"split_job_{unique_id}")
    os.makedirs(job_folder, exist_ok=True)
    try:
        reader = PdfReader(input_path)
        for i, page in enumerate(reader.pages):
            writer = PdfWriter()
            writer.add_page(page)
            with open(os.path.join(job_folder, f"page_{i+1}.pdf"), "wb") as f:
                writer.write(f)
        zip_path = os.path.join(OUTPUT_DIR, f"split_{unique_id}.zip")
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for root, dirs, files in os.walk(job_folder):
                for file in files: zipf.write(os.path.join(root, file), arcname=file)
        return zip_path
    finally:
        shutil.rmtree(job_folder, ignore_errors=True)

# NEW EDIT TOOLS: Watermark & Password Protection
def watermark_pdf_converter(input_path: str, text: str) -> str:
    output_path = os.path.join(OUTPUT_DIR, f"watermarked_{uuid.uuid4()}.pdf")
    doc = fitz.open(input_path)
    for page in doc:
        # Place a grey stamp in the middle-ish of the page
        p = fitz.Point(50, page.rect.height / 2)
        page.insert_text(p, text, fontsize=45, color=(0.5, 0.5, 0.5))
    doc.save(output_path)
    doc.close()
    return output_path

def protect_pdf_converter(input_path: str, password: str) -> str:
    output_path = os.path.join(OUTPUT_DIR, f"protected_{uuid.uuid4()}.pdf")
    doc = fitz.open(input_path)
    # Encrypts the PDF. Users will need the password to open it!
    doc.save(output_path, encryption=fitz.PDF_ENCRYPT_AES_256, owner_pw=password, user_pw=password)
    doc.close()
    return output_path

def unlock_pdf_converter(input_path: str, password: str) -> str:
    output_path = os.path.join(OUTPUT_DIR, f"unlocked_{uuid.uuid4()}.pdf")
    doc = fitz.open(input_path)
    if doc.is_encrypted:
        success = doc.authenticate(password)
        if not success:
            raise Exception("Incorrect password!")
    doc.save(output_path)
    doc.close()
    return output_path

def remove_pages_converter(input_path: str, pages_to_remove: str) -> str:
    output_path = os.path.join(OUTPUT_DIR, f"removed_{uuid.uuid4()}.pdf")
    doc = fitz.open(input_path)
    try:
        # Convert "1, 3, 5" into[0, 2, 4] for programming index
        pages =[int(p.strip()) - 1 for p in pages_to_remove.split(",") if p.strip().isdigit()]
        pages_to_keep =[i for i in range(len(doc)) if i not in pages]
        doc.select(pages_to_keep)
    except Exception:
        pass
    doc.save(output_path)
    doc.close()
    return output_path

def rotate_pdf_converter(input_path: str, degrees: str) -> str:
    output_path = os.path.join(OUTPUT_DIR, f"rotated_{uuid.uuid4()}.pdf")
    doc = fitz.open(input_path)
    try:
        deg = int(degrees)
    except:
        deg = 90
    for page in doc:
        page.set_rotation(page.rotation + deg)
    doc.save(output_path)
    doc.close()
    return output_path

# --- TO PDF CONVERTERS (Word, Excel, PPT) ---
def word_to_pdf_converter(input_path: str) -> str:
    from docx2pdf import convert
    import pythoncom
    pythoncom.CoInitialize() # Required for async web servers
    
    output_path = os.path.join(OUTPUT_DIR, f"converted_{uuid.uuid4()}.pdf")
    convert(os.path.abspath(input_path), os.path.abspath(output_path))
    return output_path

def excel_to_pdf_converter(input_path: str) -> str:
    import win32com.client
    import pythoncom
    pythoncom.CoInitialize()
    
    output_path = os.path.join(OUTPUT_DIR, f"converted_{uuid.uuid4()}.pdf")
    excel = win32com.client.Dispatch("Excel.Application")
    excel.Visible = False
    wb = excel.Workbooks.Open(os.path.abspath(input_path))
    # 0 is the code for PDF format in Excel
    wb.ExportAsFixedFormat(0, os.path.abspath(output_path))
    wb.Close(False)
    excel.Quit()
    return output_path

def ppt_to_pdf_converter(input_path: str) -> str:
    import win32com.client
    import pythoncom
    pythoncom.CoInitialize()
    
    output_path = os.path.join(OUTPUT_DIR, f"converted_{uuid.uuid4()}.pdf")
    powerpoint = win32com.client.Dispatch("Powerpoint.Application")
    # Open silently in background
    deck = powerpoint.Presentations.Open(os.path.abspath(input_path), WithWindow=False)
    # 32 is the code for PDF format in PowerPoint
    deck.SaveAs(os.path.abspath(output_path), 32)
    deck.Close()
    powerpoint.Quit()
    return output_path

# --- IMAGE TO PDF (NEW!) ---
def images_to_pdf_converter(input_paths: list) -> str:
    output_path = os.path.join(OUTPUT_DIR, f"converted_{uuid.uuid4()}.pdf")
    images =[]
    
    for p in input_paths:
        img = Image.open(p)
        # Convert images with transparency/alpha channels to RGB to avoid PDF errors
        if img.mode in ('RGBA', 'LA', 'P'):
            bg = Image.new('RGB', img.size, (255, 255, 255))
            if img.mode == 'P':
                img = img.convert('RGBA')
            bg.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
            images.append(bg)
        else:
            images.append(img.convert('RGB'))
            
    if not images:
        raise Exception("No valid images found.")
        
    # Pillow creates the PDF out of all the images seamlessly
    images[0].save(
        output_path,
        "PDF",
        resolution=100.0,
        save_all=True,
        append_images=images[1:]
    )
    return output_path

def add_page_numbers_converter(input_path: str, starting_number: str) -> str:
    output_path = os.path.join(OUTPUT_DIR, f"numbered_{uuid.uuid4()}.pdf")
    doc = fitz.open(input_path)
    
    try:
        start_num = int(starting_number)
    except:
        start_num = 1 # Default to 1 if user leaves it blank or types text
        
    for i, page in enumerate(doc):
        # Create a tiny invisible rectangle at the bottom-center of the page
        rect = fitz.Rect(0, page.rect.height - 50, page.rect.width, page.rect.height - 20)
        # 1 is the PyMuPDF code for "Center Alignment"
        page.insert_textbox(rect, str(i + start_num), fontsize=11, fontname="helv", align=1)
        
    doc.save(output_path)
    doc.close()
    return output_path

def extract_images_converter(input_path: str) -> str:
    """Extract all embedded images from a PDF into a ZIP folder"""
    unique_id = str(uuid.uuid4())
    job_folder = os.path.join(OUTPUT_DIR, f"extracted_imgs_{unique_id}")
    os.makedirs(job_folder, exist_ok=True)
    
    try:
        doc = fitz.open(input_path)
        img_count = 0
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            image_list = page.get_images(full=True)
            
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                
                img_count += 1
                image_filename = os.path.join(job_folder, f"image_{img_count}.{image_ext}")
                with open(image_filename, "wb") as f:
                    f.write(image_bytes)
                    
        doc.close()
        
        if img_count == 0:
            raise Exception("No images found inside this PDF.")
            
        zip_path = os.path.join(OUTPUT_DIR, f"extracted_images_{unique_id}.zip")
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for root, dirs, files in os.walk(job_folder):
                for file in files:
                    zipf.write(os.path.join(root, file), arcname=file)
                    
        return zip_path
    finally:
        shutil.rmtree(job_folder, ignore_errors=True)

def reorder_pages_converter(input_path: str, page_order: str) -> str:
    """Reorder pages based on user input (e.g. '3, 1, 2')"""
    output_path = os.path.join(OUTPUT_DIR, f"reordered_{uuid.uuid4()}.pdf")
    doc = fitz.open(input_path)
    
    try:
        # Convert user input "3, 1, 2" into programming indexes[2, 0, 1]
        raw_pages =[int(p.strip()) - 1 for p in page_order.split(",") if p.strip().isdigit()]
        
        # Filter out any accidental numbers that are higher than the total page count
        valid_pages =[p for p in raw_pages if 0 <= p < len(doc)]
        
        if not valid_pages:
            raise Exception("Please enter valid page numbers separated by commas.")
            
        # Reorder the PDF
        doc.select(valid_pages)
    except Exception as e:
        doc.close()
        raise Exception(f"Failed to reorder: {str(e)}. Make sure to use commas (e.g., 3,1,2)")
        
    doc.save(output_path)
    doc.close()
    return output_path


# ------------------ PDF CONVERSION API ------------------
@app.post("/api/pdf/convert")
async def convert_pdf(files: List[UploadFile] = File(...), tool: str = Form(...), extra_param: Optional[str] = Form("")):
    input_paths =[]
    for f in files:
        path = os.path.join(UPLOAD_DIR, f"{uuid.uuid4()}_{f.filename}")
        with open(path, "wb") as buffer:
            shutil.copyfileobj(f.file, buffer)
        input_paths.append(path)

    # For non-merge tools, just use the first file
    input_path = input_paths[0]
    
    try:
        output_path = None
        filename = "converted"
        
        if tool == "pdf-to-word":
            output_path = pdf_to_word_converter(input_path)
            filename = "converted.docx"

        elif tool == "pdf-to-text":
            output_path = pdf_to_text_converter(input_path)
            filename = "converted.txt"

        elif tool == "pdf-to-mp3":
            output_path = pdf_to_mp3_converter(input_path)
            filename = "converted.mp3"

        elif tool == "pdf-to-excel":
            output_path = pdf_to_excel_converter(input_path)
            filename = "converted.xlsx"

        elif tool == "pdf-to-ppt":
            output_path = pdf_to_ppt_converter(input_path)
            filename = "converted.pptx"

        elif tool == "pdf-to-image":
            output_path = pdf_to_image_converter(input_path)
            filename = "images.zip"

        elif tool == "compress-pdf":
            output_path = compress_pdf_converter(input_path)
            filename = "compressed.pdf"

        elif tool == "merge-pdf":
            if len(input_paths) < 2:
                return JSONResponse(status_code=400, content={"error": "Please upload at least 2 PDF files to merge."})
            output_path = merge_pdf_converter(input_paths)
            filename = "merged.pdf"
            
        elif tool == "image-to-pdf":
            if len(input_paths) < 1:
                return JSONResponse(status_code=400, content={"error": "Please upload at least 1 image."})
            output_path = images_to_pdf_converter(input_paths)
            filename = "images_converted.pdf"

        elif tool == "split-pdf":
            output_path = split_pdf_converter(input_path)
            filename = "split_pages.zip"

        elif tool == "word-to-pdf":
            output_path = word_to_pdf_converter(input_path)
            filename = "document.pdf"

        elif tool == "excel-to-pdf":
            output_path = excel_to_pdf_converter(input_path)
            filename = "spreadsheet.pdf"

        elif tool == "ppt-to-pdf":
            output_path = ppt_to_pdf_converter(input_path)
            filename = "presentation.pdf"

        elif tool == "watermark-pdf":
            text_to_stamp = extra_param if extra_param else "Confidential"
            output_path = watermark_pdf_converter(input_path, text_to_stamp)
            filename = "watermarked.pdf"

        elif tool == "protect-pdf":
            password_to_set = extra_param if extra_param else "password"
            output_path = protect_pdf_converter(input_path, password_to_set)
            filename = "protected_file.pdf"

        elif tool == "unlock-pdf":
            output_path = unlock_pdf_converter(input_path, extra_param or "")
            filename = "unlocked.pdf"

        elif tool == "remove-pages":
            output_path = remove_pages_converter(input_path, extra_param or "1")
            filename = "cleaned.pdf"
            
        elif tool == "rotate-pdf":
            output_path = rotate_pdf_converter(input_path, extra_param or "90")
            filename = "rotated.pdf"

        elif tool == "add-page-numbers":
            output_path = add_page_numbers_converter(input_path, extra_param or "1")
            filename = "numbered.pdf"

        elif tool == "extract-images":
            output_path = extract_images_converter(input_path)
            filename = "extracted_images.zip"

        elif tool == "reorder-pages":
            if not extra_param:
                return JSONResponse(status_code=400, content={"error": "Please provide the new page order."})
            output_path = reorder_pages_converter(input_path, extra_param)
            filename = "reordered.pdf"

        else:
            return JSONResponse(status_code=400, content={"error": f"Unsupported tool: {tool}"})
        
        if not output_path or not os.path.exists(output_path):
            raise Exception("File processing failed - No output generated")
        
        # Give the frontend a SECURE LINK instead of the raw file
        file_id = os.path.basename(output_path)
        return JSONResponse({"success": True, "download_url": f"/api/download/{file_id}?name={filename}"})
    
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})
    
    finally:
        for p in input_paths:
            if os.path.exists(p):
                os.remove(p)

# SECURE DOWNLOAD ENDPOINT (Defeats Download Managers)
@app.get("/api/download/{file_id}")
async def download_file(file_id: str, name: str = "download", background_tasks: BackgroundTasks = BackgroundTasks()):
    if ".." in file_id or "/" in file_id or "\\" in file_id:
        return JSONResponse(status_code=400, content={"error": "Invalid file request"})
        
    file_path = os.path.join(OUTPUT_DIR, file_id)
    if not os.path.exists(file_path):
        return JSONResponse(status_code=404, content={"error": "File already downloaded or expired."})
    
    # Wait 300 seconds (5 minutes) before deleting so IDM can finish downloading!
    background_tasks.add_task(delete_file_delayed, file_path, 300)
    
    return FileResponse(file_path, media_type="application/octet-stream", filename=name)

# ------------------ IMAGE TOOLS ------------------
app.include_router(tools_router)

# ------------------ RUN ------------------
if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000, reload=True)